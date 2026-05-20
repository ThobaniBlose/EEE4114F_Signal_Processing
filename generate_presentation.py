"""
Generate SHARC Group 23 project presentation.
Run: pip install python-pptx && python generate_presentation.py
Output: SHARC_Group23_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helpers
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(title, bullets, note="", diagram_note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
            tf.paragraphs[0].font.size = Pt(18)
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.space_before = Pt(6)
    if diagram_note:
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = f"[{diagram_note}]"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

# --- SLIDE 0: COVER ---
add_title_slide(
    "A Low-Cost Autonomous Buoy for\nWave-Direction Measurement in the\nAntarctic Marginal Ice Zone",
    "EEE4113F 2026 — Group 23\n\n"
    "Batsirai Rwatirera — Sensing\n"
    "Karen Karenyi — Power\n"
    "Thobani Blose — Signal Processing\n"
    "Nyeleti Mushwana — Housing\n\n"
    "Department of Electrical Engineering\nUniversity of Cape Town"
)

# --- SLIDE 1: INTRODUCTION ---
add_content_slide(
    "Introduction — Why Directional Wave Measurement Matters",
    [
        "The Antarctic marginal ice zone needs better wave monitoring for understanding wave–ice interaction and ocean energy budgets",
        "Current low-cost polar sensors measure only heave/frequency spectra — they show energy magnitude but not propagation direction",
        "Polar deployment is uniquely difficult: ice, extreme cold, limited access, equipment loss during ice-melt, restricted communication",
        "Project purpose: design a low-cost autonomous buoy that measures wave direction, not just wave frequency",
    ],
    note="Emphasise that directional information separates this project from existing low-cost heave-only sensors.",
    diagram_note="DIAGRAM: Map of Antarctic MIZ or polar wave environment photo"
)

# --- SLIDE 2: PROBLEM STATEMENT ---
add_content_slide(
    "Problem Statement and User Need",
    [
        "Non-directional sensors discard the wave propagation vector — researchers cannot distinguish wave direction",
        "Engineering challenges: survivability, waterproofing, onboard processing, compact output, unreliable magnetic heading at high latitudes",
        "Stakeholder: Ms Robyn Verrinder (polar marine researcher, UCT) — needs compact directional data from remote deployments",
        "D-school process refined the problem from 'ice-melt survival' to 'autonomous wave-direction measurement'",
    ],
    note="The problem was refined through D-school activities including empathy mapping, expert interviews, and iterative prototyping.",
    diagram_note="DIAGRAM: Empathy map or stakeholder needs summary"
)

# --- SLIDE 3: SYSTEM ARCHITECTURE ---
add_content_slide(
    "Proposed System Architecture",
    [
        "Sensing: IMU + GPS/GNSS → time-stamped motion data → SD card",
        "Signal Processing: raw motion → filtering → integration → spectral estimation → wave parameters + direction",
        "Power: battery → regulated 5V/3.3V → switching → monitoring → protection",
        "Housing: waterproof enclosure → buoyancy → thermal protection → structural integrity",
        "Data flow: IMU/GPS → CSV storage → onboard processing → compact Tier-1 output packet",
    ],
    note="The system is self-contained and autonomous. Each subsystem was designed and tested independently with defined interfaces.",
    diagram_note="DIAGRAM: System block diagram showing four subsystems and data flow"
)

# --- SLIDE 4: SENSING ---
add_content_slide(
    "Sensing Subsystem — Data Acquisition",
    [
        "Acquires 3-axis accel/gyro/mag (MPU-6050, 100 Hz) and GPS/GNSS (8 Hz)",
        "Stores time-stamped 13-column CSV on SD card (S001_IMU.csv format)",
        "Validated: sampling rate, GPS update, CSV format, timestamp sync, SD write path",
        "Limitations: GPS antenna reliability, mid-session fault handling, magnetometer degradation at high latitudes",
    ],
    note="The sensing subsystem defines the data interface that all downstream processing depends on.",
    diagram_note="DIAGRAM: Sensing block diagram or CSV format example"
)

# --- SLIDE 5: POWER ---
add_content_slide(
    "Power Subsystem — Regulated Supply and Protection",
    [
        "Li-SOCl₂ battery → buck converter → 5V and 3.3V regulated rails",
        "High-side switching for GPS/SD power, reverse-polarity protection, battery monitoring",
        "ATPs passed: 1-hour endurance, ripple < 50 mVpp, switched rails, reverse polarity protection",
        "Known limitation: 5V rail requires ≥5.8V input (original spec was 5.4V)",
        "Recommendations: PTC fuse, improved gate-drive margin, updated battery threshold",
    ],
    note="The voltage headroom issue is a design margin problem resolvable with component selection changes.",
    diagram_note="DIAGRAM: Power subsystem schematic or block diagram"
)

# --- SLIDE 6: SIGNAL PROCESSING ---
add_content_slide(
    "Signal-Processing Subsystem — Motion to Wave Parameters",
    [
        "Method: MATLAB golden reference → validated STM32 C implementation",
        "Pipeline: az[g] → filtfilt → integrate → η(t) → Welch PSD → Hm0, Tp, Tm01, Tm02 → cross-spectral direction",
        "STM32 reproduced MATLAB with 0.002% Hm0 error, <0.001% period errors",
        "Automatic BODY_RELATIVE fallback when heading/magnetometer unavailable",
        "ATP validation: geographic direction verified with <0.001° error (synthetic cases)",
        "Limitation: full 30-min SD-card processing not demonstrated (RAM-limited to 5.5-min segment)",
    ],
    note="The signal-processing subsystem is the most thoroughly validated component. Every stage verified against MATLAB to machine precision.",
    diagram_note="DIAGRAM: Signal processing pipeline or validation results table"
)

# --- SLIDE 7: HOUSING ---
add_content_slide(
    "Housing Subsystem — Protection and Buoyancy",
    [
        "Cylindrical enclosure with internal board-mounting cradle",
        "Pressure equalisation, O-ring sealing, magnetic activation switch",
        "Analysis: buoyancy/orientation, thermal simulation, structural loading",
        "PLA prototype manufactured; HDPE recommended for deployment",
        "Recommendations: IP67 submersion test, thermal cycling, saltwater/UV exposure testing",
    ],
    note="The housing design is analytically justified but not validated under real deployment conditions.",
    diagram_note="DIAGRAM: Housing CAD render, prototype photo, or simulation results"
)

# --- SLIDE 8: CONCLUSION ---
add_content_slide(
    "Conclusion and Future Work",
    [
        "Proof-of-concept achieved: four subsystems individually designed and validated",
        "Signal processing matched MATLAB reference to <0.01% across all wave parameters",
        "Automatic mode decision (GEOGRAPHIC / BODY_RELATIVE / FALLBACK) proven from data",
        "This is a laboratory prototype — not a deployment-ready polar instrument",
        "Future: full integration test, 30-min SD processing, heading-valid field data, wave-tank validation with independent reference",
    ],
    note="The project achieved its proof-of-concept objective. The path to a deployable instrument is clear but requires integration testing and environmental validation.",
    diagram_note="DIAGRAM: Summary table or development roadmap"
)

# Save
output_path = "SHARC_Group23_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
